"""
Document to PDF converter module.

Supports converting various document formats (PPT, PPTX) to PDF.
"""
from __future__ import annotations

import os
import tempfile
from pathlib import Path
from typing import Literal

SUPPORTED_EXTENSIONS = {".ppt", ".pptx", ".pdf"}

ConversionResult = Literal["converted", "already_pdf", "unsupported"]


class DocumentConversionError(Exception):
    """Raised when document conversion fails."""


def is_supported_format(file_path: str) -> bool:
    """Check if the file format is supported for material import."""
    ext = Path(file_path).suffix.lower()
    return ext in SUPPORTED_EXTENSIONS


def convert_to_pdf(source_path: str, output_path: str) -> ConversionResult:
    """
    Convert a document to PDF format.
    
    Args:
        source_path: Path to the source document
        output_path: Path where the PDF should be saved
        
    Returns:
        "converted" if conversion was performed
        "already_pdf" if source was already a PDF (copied directly)
        "unsupported" if file format is not supported
        
    Raises:
        DocumentConversionError: If conversion fails
        FileNotFoundError: If source file doesn't exist
    """
    source = Path(source_path)
    if not source.exists():
        raise FileNotFoundError(f"Source file not found: {source_path}")
    
    ext = source.suffix.lower()
    
    if ext == ".pdf":
        # Already PDF, just copy it
        import shutil
        shutil.copy2(source_path, output_path)
        return "already_pdf"
    
    if ext not in SUPPORTED_EXTENSIONS:
        return "unsupported"
    
    if ext in {".ppt", ".pptx"}:
        return _convert_pptx_to_pdf(source_path, output_path)
    
    return "unsupported"


def _convert_pptx_to_pdf(pptx_path: str, pdf_path: str) -> ConversionResult:
    """
    Convert PPT/PPTX to PDF using python-pptx and reportlab.
    
    This is a simplified conversion that extracts slides as images
    and creates a PDF. For better results, consider using LibreOffice
    or other dedicated converters.
    """
    try:
        from pptx import Presentation
        from reportlab.lib.pagesizes import letter
        from reportlab.pdfgen import canvas
        from PIL import Image
        import io
    except ImportError as exc:
        raise DocumentConversionError(
            "Required libraries not installed. Run: pip install python-pptx reportlab pillow"
        ) from exc
    
    try:
        # Load the presentation
        prs = Presentation(pptx_path)
        
        if not prs.slides:
            raise DocumentConversionError("Presentation has no slides")
        
        # Create PDF
        c = canvas.Canvas(pdf_path, pagesize=letter)
        width, height = letter
        
        # For simplicity, we'll create a text-based PDF with slide content
        # A full implementation would render slides as images
        for idx, slide in enumerate(prs.slides, 1):
            # Extract text from slide
            text_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text_content.append(shape.text)
            
            # Write to PDF
            c.setFont("Helvetica-Bold", 16)
            c.drawString(50, height - 50, f"Slide {idx}")
            
            c.setFont("Helvetica", 12)
            y_position = height - 100
            for text in text_content:
                # Simple text wrapping
                lines = _wrap_text(text, 80)
                for line in lines:
                    if y_position < 50:
                        c.showPage()
                        c.setFont("Helvetica", 12)
                        y_position = height - 50
                    c.drawString(50, y_position, line[:100])
                    y_position -= 20
            
            # New page for next slide
            if idx < len(prs.slides):
                c.showPage()
        
        c.save()
        return "converted"
        
    except Exception as exc:
        raise DocumentConversionError(f"Failed to convert PPTX to PDF: {str(exc)}") from exc


def _wrap_text(text: str, width: int) -> list[str]:
    """Simple text wrapping helper."""
    words = text.split()
    lines = []
    current_line = []
    current_length = 0
    
    for word in words:
        word_length = len(word) + 1
        if current_length + word_length > width:
            if current_line:
                lines.append(" ".join(current_line))
                current_line = [word]
                current_length = word_length
            else:
                lines.append(word)
                current_length = 0
        else:
            current_line.append(word)
            current_length += word_length
    
    if current_line:
        lines.append(" ".join(current_line))
    
    return lines or [""]
